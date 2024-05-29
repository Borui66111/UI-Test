from pptx import Presentation
from pptx.util import Pt,Inches,Emu
from collections import defaultdict as dd
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import numpy as np, matplotlib.pyplot as plt, cv2, os

'''
*Notes*

* slide = prs.slides.add_slide(slide_layout) adds a new slide with predefined layout, however, how to add one text chunk alone?
* prs.slides[slide].shapes[ele].add_picture(img,location) adds a new picture 
* add text chunk
  * prs.slides[slide].shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
  * txBox = slide.shapes.add_textbox(left, top, width, height); tf = txBox.text_frame; tf.text = item["text"]
* layouts (0-8):
  * Title (presentation title slide)
  * Title and Content
  * Section Header (sometimes called Segue)
  * Two Content (side by side bullet textboxes)
  * Comparison (same but additional title for each side by side content box)
  * Title Only
  * Blank
  * Content with Caption
  * Picture with Caption
* format
  * text_frame.vertical_anchor = MSO_ANCHOR.TOP
  * text_frame.word_wrap = False
  * text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

* These coordinates are defined using the standard units of points in PowerPoint (1 inch = 72 point = 96 pixel = 914400 emu (English Metric Unit)). The top-left corner of the slide is considered the origin point (0,0). A common ppt slide is 10x7.5 inches(4:3,default) or 13.3x7.5 inches(16:9).
* location: [w,h]

ref: https://python-pptx.readthedocs.io/en/latest/api/text.html#pptx.text.text.TextFrame
'''

path='/Users/Borui/Dean/Computing/AI/Others/UROP/PPT_Generation/Self Introduction - Interns Onboarding.pptx'

class PPT:
    prs=Presentation()
    prs.slides.add_slide(prs.slide_layouts[0])
    shape=[Inches(10),Inches(7.5)]
    def __init__(self,prs='') -> None:
        self.prs=Presentation(prs) if prs else Presentation()
    def to_emu(self,bbox:[]):
        return [Inches(i) for i in bbox] if bbox else None
    def add_slide(self,prs:Presentation=None,slide_layout=0):
        '''return slide'''
        prs=prs if prs else self.prs;slide_layout = prs.slide_layouts[slide_layout]
        slide = prs.slides.add_slide(slide_layout);return slide
    # def pxl_to_emu(self,bbox=['x','y',None,None],shape=[1024,768]):
    #     ratio=PPT.shape[0].inches/shape[0]; return [i*ratio if i else i for i in bbox]
    # def inch_to_pxl(self,bbox=['x Inches','y Inches',None,None],shape=[1024,768]):
    #     ratio=shape[0]/PPT.shape[0].inches; return [i.inches*ratio if i else i for i in bbox]
    def add_text_chunk(self,shapes=prs.slides[0].shapes,bbox=[Inches(0),Inches(0),shape[0]/4,shape[1]/4]):
        '''return shape (text frame)'''
        shape=shapes.add_textbox(*[x if type(x) in [Inches,Emu] else Emu(x) for x in bbox])
        shape.text_frame.word_wrap=True;return shape
    def get_shape(self,prs:Presentation=prs,slide=0,ele=0):
        prs=prs if prs else self.prs; return prs.slides[slide].shapes[ele]
    def change_text(self,shape,text='Hello World!'):
        shape.text=text
    def set_font(self,shape,name='Arial',size=Pt(18),color=(0,0,0),bold=None,italic=None,underline=None,word_wrap=True):
        for paragraph in (shape.text_frame.paragraphs):
            font=paragraph.font;font.name = name;font.size=Pt(size);font.color.rgb = RGBColor(*color);font.bold=bold;font.italic=italic;font.underline=underline
        shape.text_frame.word_wrap=word_wrap
    def get_pos(self,shape:prs.slides[0].shapes[0]):
        return [shape.left,shape.top,shape.width,shape.height]
    def set_pos(self,shape:prs.slides[0].shapes[0],bbox=to_emu([0,0],None)):
        '''bbox: unit=emu'''
        shape.left,shape.top,shape.width,shape.height=bbox if len(bbox)==4 else bbox+[shape.width,shape.height]
    def add_img(self,slide=prs.slides[0],pos=[Inches(0),Inches(0)],img='sample.jpg',width=None,height=None):
        slide.shapes.add_picture(img,*[x if type(x) in [Inches,Emu,type(None)] else Inches(x) for x in pos+[width,height]])
    def prs_save(self,prs:Presentation=None,name='../ppt/test.pptx'):
        prs=prs if prs else self.prs; prs.save(name)
    def prs_show_info(self,prs=None):
        prs=prs if prs else self.prs;print(f"PPT: width={prs.slide_width}; height={prs.slide_height}")
        for slide_number, slide in enumerate(prs.slides, start=1):
            print(f"Slide {slide_number} Elements:")
            for shape_number, shape in enumerate(slide.shapes, start=1):
                x = shape.left; y = shape.top; w = shape.width; h = shape.height; print(f"Element {shape_number}:\nshape: {shape}, x: {x}, y: {y}, width: {w}, height: {h}")
    def get_img_info(self,img):
        img = img if type(img)==np.ndarray else cv2.imread(img);print(img.shape)
        cv2.imshow('img',img);cv2.waitKey(-1);cv2.destroyAllWindows();return img
    def load_imgs(self,folder='./images/',shape=(0,0),fx=0.3,fy=0.3,display=False):
        for im in [os.path.join(folder,ii) for ii in os.listdir(folder)]:
            im=cv2.resize(cv2.imread(im),shape,fx=fx,fy=fy)
            if display: cv2.imshow('',im);cv2.waitKey(-1);cv2.destroyAllWindows(); 
            yield im
    

